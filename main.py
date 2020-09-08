import math, operator

import pptx
from PIL import Image
import sys
import os
import glob
import subprocess
import shutil
from functools import reduce

from pptx import Presentation


def get_page(url):
    import urllib3
    http = urllib3.PoolManager()
    r = http.request('GET', url)
    print((r.status))
    print((r.data))

    return None

# this was found on some stackoverflow
def are_image_same(file1, file2):
    image1 = Image.open(file1)
    image2 = Image.open(file2)

    h1 = image1.histogram()
    h2 = image2.histogram()
    rms = math.sqrt(reduce(operator.add, list(map(lambda a,b: (a-b)**2, h1, h2)))/len(h1))

    return rms == 0


def to_power_point(img_folder):
    prs = Presentation()

    filelist = glob.glob(os.path.join(img_folder, '*.png'))
    filelist.sort()
    for file in filelist:
        title_slide_layout = prs.slide_layouts[6] # 6 is blank
        slide = prs.slides.add_slide(title_slide_layout)
        width = pptx.util.Cm(33.86)
        slide.shapes.add_picture(file, 0, 0, width=width)

    prs.save('./test.pptx')

# Press the green button in the gutter to run the script.
if __name__ == '__main__':
    url = 'https://sitescours.monportail.ulaval.ca/ena/site/module?editionModule=false&idSite=121706&idPage=2790239&idModule=1030262&_js=true'
    # get_page(url)

    video = "C:/Users/naked_000/Desktop/video(1).mp4"
    out_dir = "./out"
    if not os.path.exists("decomp"):
        os.mkdir("decomp")
    else:
        sys.exit("decomp already exists, exit")

    if not os.path.exists(out_dir):
        os.mkdir(out_dir)

    cmd = ["ffmpeg", "-i", video, "-vf", "select='eq(pict_type,I)'", "-vsync", "0", "-f", "image2",
           "decomp/%09d.png"]

    print(("Running ffmpeg: " + " ".join(cmd)))


    subprocess.call(cmd)

    print("Done, now eliminating duplicate images and moving unique ones to output folder...")

    filelist = glob.glob(os.path.join("decomp", '*.png'))
    filelist.sort()
    for ii in range(0, len(filelist)):
        if ii < len(filelist) - 1:
            if are_image_same(filelist[ii], filelist[ii + 1]):
                print(('Found similar images: ' + filelist[ii] + " and " + filelist[ii + 1]))
            else:
                print(('Found unique image: ' + filelist[ii]))
                head, tail = os.path.split(filelist[ii])
                shutil.copyfile(filelist[ii], out_dir + os.path.sep + tail)
        else:
            shutil.copyfile(filelist[ii], out_dir + os.path.sep + tail)
    shutil.rmtree("decomp")

    to_power_point(out_dir)
