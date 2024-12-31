"""
This module generates PowerPoint presentations from image frames.
It adds each image in the specified folder as a slide in a new PowerPoint presentation,
with options to customize slide layout and output naming.
"""

import os
import glob
from pptx import Presentation
import logging

def to_power_point(img_folder, base_name, export_folder="./presentation/", slide_layout=6, output_name=None):
    """
    Converts images in a folder to a PowerPoint presentation.

    Args:
        img_folder (str): Path to the folder containing images.
        base_name (str): Base name for the output PowerPoint file.
        export_folder (str): Directory to save the PowerPoint file. Defaults to './presentation/'.
        slide_layout (int): Slide layout index to use for the presentation. Defaults to 6 (blank layout).
        output_name (str, optional): Custom name for the PowerPoint file. Defaults to None.

    Raises:
        FileNotFoundError: If the specified image folder does not exist.
        ValueError: If no PNG files are found in the specified folder.
    """
    if not os.path.exists(img_folder):
        raise FileNotFoundError(f"Image folder not found: {img_folder}")

    filelist = sorted(glob.glob(os.path.join(img_folder, '*.png')))
    if not filelist:
        raise ValueError(f"No PNG files found in {img_folder}")

    if not os.path.exists(export_folder):
        os.makedirs(export_folder)

    logging.basicConfig(level=logging.INFO)
    prs = Presentation()

    for file in filelist:
        logging.info(f"Adding image {file} to slide")
        slide = prs.slides.add_slide(prs.slide_layouts[slide_layout])
        slide_width = prs.slide_width
        slide.shapes.add_picture(file, 0, 0, width=slide_width)

    output_name = output_name or f"{base_name}.pptx"
    prs.save(os.path.join(export_folder, output_name))
    logging.info(f"Presentation saved at {os.path.join(export_folder, output_name)}")
